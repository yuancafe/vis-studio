#!/usr/bin/env python3
"""Run VIS STUDIO orchestration in document or production mode."""

from __future__ import annotations

import argparse
import json
import select
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from build_execution_plan import build_plan
from build_pipeline import build_pipeline
from build_tool_execution_payload import build_payload
from _orchestration_utils import (
    choose_executor_for_intent,
    discover_design_services,
    load_route_rules,
    run_cmd,
    slugify,
    utc_timestamp,
)

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    HAS_REPORTLAB = True
except Exception:
    HAS_REPORTLAB = False

try:
    from pptx import Presentation

    HAS_PPTX = True
except Exception:
    HAS_PPTX = False


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("goal", help="Production goal text")
    parser.add_argument(
        "--intent",
        default=None,
        help="Optional explicit intent",
    )
    parser.add_argument(
        "--run-mode",
        default="document",
        choices=["document", "production"],
        help="Run mode",
    )
    parser.add_argument(
        "--execution-mode",
        default="handoff",
        choices=["handoff", "execute"],
        help="Execution mode used in document mode",
    )
    parser.add_argument(
        "--production-scope",
        default="single",
        choices=["single", "full_suite"],
        help="Production scope",
    )
    parser.add_argument(
        "--execution-gateway-mode",
        default="hybrid",
        choices=["direct", "swarm", "hybrid"],
        help="Gateway mode for production execution.",
    )
    parser.add_argument(
        "--brand-pack",
        type=Path,
        default=None,
        help="Brand pack path",
    )
    parser.add_argument(
        "--order-file",
        type=Path,
        default=None,
        help="Optional order form path. If omitted, default order is auto-generated.",
    )
    parser.add_argument(
        "--resume-from-run-dir",
        type=Path,
        default=None,
        help="Optional previous production run directory for incremental orders.",
    )
    parser.add_argument(
        "--preferred-tool",
        default=None,
        help="Optional preferred tool",
    )
    parser.add_argument(
        "--route-rules",
        type=Path,
        default=None,
        help="Optional route rules path",
    )
    parser.add_argument(
        "--output-root",
        type=Path,
        default=Path.cwd(),
        help="Output root path (default: current directory)",
    )
    parser.add_argument(
        "--decision-timeout-seconds",
        type=int,
        default=15,
        help="Timeout in seconds for critical-step decisions",
    )
    parser.add_argument(
        "--fallback-figma-connector",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Allow figma connector fallback when figma MCP is unavailable",
    )
    parser.add_argument(
        "--executor-mode",
        choices=["probe", "mock"],
        default="probe",
        help="Executor behavior for production mode",
    )
    return parser.parse_args()


def _write_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    text = content if content.endswith("\n") else content + "\n"
    path.write_text(text, encoding="utf-8")


def _normalize_order_form(order: dict[str, Any], resume_override: Path | None = None) -> dict[str, Any]:
    selected = order.get("selected_deliverables", [])
    if not isinstance(selected, list):
        selected = []
    selected = [str(item).strip() for item in selected if str(item).strip()]
    if not selected:
        selected = ["logo_core_package"]

    manual_formats = order.get("manual_formats", [])
    if not isinstance(manual_formats, list):
        manual_formats = []
    normalized_formats = []
    for item in manual_formats:
        fmt = str(item).strip().lower()
        if fmt in {"pdf", "ai", "pptx"}:
            normalized_formats.append(fmt)
    if not normalized_formats:
        normalized_formats = ["pdf", "ai"]

    resume_from = ""
    if resume_override is not None:
        resume_from = str(resume_override)
    elif isinstance(order.get("resume_from_run_dir"), str):
        resume_from = order.get("resume_from_run_dir", "").strip()

    constraints = order.get("client_constraints", {})
    if not isinstance(constraints, dict):
        constraints = {}

    return {
        "selected_deliverables": list(dict.fromkeys(selected)),
        "manual_formats": list(dict.fromkeys(normalized_formats)),
        "auto_generate_application_assets": bool(order.get("auto_generate_application_assets", False)),
        "resume_from_run_dir": resume_from,
        "client_constraints": constraints,
    }


def _load_or_create_order_form(
    order_file: Path | None,
    run_root: Path,
    resume_override: Path | None,
) -> tuple[dict[str, Any], Path]:
    default_order = {
        "selected_deliverables": ["logo_core_package"],
        "manual_formats": ["pdf", "ai"],
        "auto_generate_application_assets": False,
        "resume_from_run_dir": str(resume_override) if resume_override else "",
        "client_constraints": {},
    }
    if order_file is not None and order_file.exists():
        try:
            loaded = json.loads(order_file.read_text(encoding="utf-8"))
            if isinstance(loaded, dict):
                normalized = _normalize_order_form(loaded, resume_override=resume_override)
            else:
                normalized = default_order
        except Exception:
            normalized = default_order
    else:
        normalized = default_order

    output_order_file = run_root / "ORDER_FORM.json"
    _write_json(output_order_file, normalized)
    return normalized, output_order_file


def _build_setup_required_doc(path: Path, discovery_report: dict[str, Any]) -> None:
    lines = [
        "# MCP Setup Required",
        "",
        "The production run detected missing critical MCP execution capability.",
        "",
        "## Missing critical routes",
    ]
    missing = discovery_report.get("missing_critical", [])
    if not missing:
        lines.append("- none")
    else:
        for item in missing:
            lines.append(
                f"- `{item.get('intent')}` requires one of: {', '.join(item.get('executor_candidates', []))}"
            )
    lines.extend(["", "## Suggested config snippets"])
    suggestions = discovery_report.get("suggestions", [])
    if not suggestions:
        lines.append("- no suggestions generated")
    else:
        for suggestion in suggestions:
            lines.append(f"### {suggestion.get('intent')}")
            lines.append("")
            lines.append("```toml")
            lines.append(suggestion.get("suggested_config", "<add server config here>"))
            lines.append("```")
            lines.append("")
    lines.extend(
        [
            "## Next-step commands",
            "",
            "```bash",
            "python3 support/brand-production-orchestrator/scripts/discover_mcp_servers.py --output MCP_DISCOVERY_REPORT.json",
            "```",
        ]
    )
    _write_text(path, "\n".join(lines))


def _build_app_cli_setup_doc(path: Path, missing_apps: list[str]) -> None:
    lines = [
        "# App CLI Setup Required",
        "",
        "No direct/swarm gateway was available for some routes and required local apps were not found.",
        "",
        "## Missing local apps",
    ]
    if not missing_apps:
        lines.append("- none")
    else:
        for app in missing_apps:
            lines.append(f"- {app}")
    lines.extend(
        [
            "",
            "## Recommendations",
            "- Install the missing app(s) and re-run production mode.",
            "- Or configure corresponding MCP server in `~/.codex/config.toml`.",
            "- Or continue in document mode and complete manually.",
        ]
    )
    _write_text(path, "\n".join(lines))


def _prompt_for_critical_decision(task_id: str, timeout_seconds: int) -> str:
    print(
        f"[critical-failure] {task_id}: choose retry/skip/stop (default stop in {timeout_seconds}s):",
        flush=True,
    )
    readable, _, _ = select.select([sys.stdin], [], [], max(1, timeout_seconds))
    if not readable:
        return "stop"
    response = sys.stdin.readline().strip().lower()
    if response in {"retry", "skip", "stop"}:
        return response
    return "stop"


def _select_executor_for_task(
    task: dict[str, Any],
    discovery_report: dict[str, Any],
    fallback_figma_connector: bool,
) -> tuple[str, list[dict[str, Any]], list[dict[str, str]], str, dict[str, Any]]:
    selected, candidates, selected_reason = choose_executor_for_intent(task["intent"], discovery_report)
    filtered_candidates: list[dict[str, Any]] = []
    for candidate in candidates:
        item = dict(candidate)
        if item.get("executor") == "figma_connector" and not fallback_figma_connector:
            item["available"] = False
            item["reason"] = "figma-connector-fallback-disabled"
        filtered_candidates.append(item)

    selected_candidate: dict[str, Any] = {}
    if selected:
        selected_candidate = next(
            (item for item in filtered_candidates if item.get("executor") == selected),
            {},
        )
    if not selected_candidate:
        selected_candidate = next((item for item in filtered_candidates if item.get("available")), {})
        selected = selected_candidate.get("executor", "")
        if selected_candidate:
            selected_reason = "selected-from-priority-candidates"

    fallback_trace: list[dict[str, str]] = []
    preferred = task.get("executor_candidates", [])
    if preferred and selected and preferred[0] != selected:
        fallback_trace.append(
            {
                "from": preferred[0],
                "to": selected,
                "reason": "higher-priority executor unavailable",
            }
        )
    return selected, filtered_candidates, fallback_trace, selected_reason, selected_candidate


def _collect_dependency_inputs(task: dict[str, Any], artifact_registry: dict[str, Any]) -> list[dict[str, Any]]:
    deps: list[dict[str, Any]] = []
    tasks = artifact_registry.get("tasks", {})
    allowed_types = set(task.get("input_artifact_types", []))
    for dep_task_id in task.get("depends_on", []):
        dep = tasks.get(dep_task_id, {})
        dep_added = 0
        for item in dep.get("output_files", []):
            artifact_type = str(item.get("artifact_type", ""))
            if allowed_types and artifact_type and artifact_type not in allowed_types:
                continue
            deps.append(
                {
                    "task_id": dep_task_id,
                    "package": dep.get("package", ""),
                    "artifact_type": artifact_type,
                    "format": item.get("format", ""),
                    "path": item.get("path", ""),
                }
            )
            dep_added += 1
        if dep_added == 0:
            for item in dep.get("output_files", []):
                deps.append(
                    {
                        "task_id": dep_task_id,
                        "package": dep.get("package", ""),
                        "artifact_type": str(item.get("artifact_type", "")),
                        "format": item.get("format", ""),
                        "path": item.get("path", ""),
                    }
                )
    return deps


def _attempt_direct_gateway(selected_candidate: dict[str, Any]) -> tuple[bool, dict[str, Any]]:
    matched = selected_candidate.get("matched_servers", [])
    if not matched:
        return False, {"reason": "direct-no-matched-server"}
    server_id = str(matched[0].get("server_id", "")).strip()
    if not server_id:
        return False, {"reason": "direct-empty-server-id"}

    return_code, stdout, stderr = run_cmd(["mcporter", "list", server_id, "--output", "json"], timeout=30)
    stderr_lower = (stderr or "").lower()
    if return_code != 0:
        if "unknown mcp server" in stderr_lower:
            return False, {"reason": "direct-unknown-server", "server_id": server_id, "stderr": stderr}
        return False, {"reason": "direct-probe-failed", "server_id": server_id, "stderr": stderr}
    try:
        json.loads(stdout or "{}")
    except json.JSONDecodeError:
        return False, {"reason": "direct-invalid-json", "server_id": server_id}
    return True, {"reason": "direct-probe-ok", "server_id": server_id}


def _attempt_swarm_gateway(discovery_report: dict[str, Any]) -> tuple[bool, dict[str, Any]]:
    env = discovery_report.get("gateway_environment", {})
    if not env.get("mcporter_available"):
        return False, {"reason": "swarm-mcporter-unavailable"}
    if not env.get("swarm_detected"):
        return False, {"reason": "swarm-not-detected"}
    # Swarm compatibility: use reachable config + execution attempt path rather than strict direct server checks.
    return True, {"reason": "swarm-managed-execution-attempted"}


def _attempt_app_cli_gateway(selected_candidate: dict[str, Any]) -> tuple[bool, dict[str, Any]]:
    cli_hits = selected_candidate.get("cli_hits", [])
    cli_probe_errors: list[dict[str, Any]] = []
    for hit in cli_hits:
        probe_command = [str(item) for item in hit.get("probe_command", []) if str(item).strip()]
        if not probe_command:
            continue
        return_code, stdout, stderr = run_cmd(probe_command, timeout=30)
        if return_code == 0:
            return True, {
                "reason": "app-cli-cli-probe-ok",
                "cli_id": str(hit.get("cli_id", "")),
                "cli_path": str(hit.get("path", "")),
                "probe_command": probe_command,
                "cli_command": [str(item) for item in hit.get("execute_command", []) if str(item).strip()],
                "stdout_preview": (stdout or "")[:240],
            }
        cli_probe_errors.append(
            {
                "cli_id": str(hit.get("cli_id", "")),
                "probe_command": probe_command,
                "return_code": return_code,
                "stderr_preview": (stderr or "")[:240],
            }
        )

    hits = selected_candidate.get("app_cli_hits", [])
    if not hits:
        if cli_probe_errors:
            return False, {
                "reason": "app-cli-cli-probe-failed",
                "cli_probe_errors": cli_probe_errors,
            }
        return False, {"reason": "app-cli-not-installed"}
    path = hits[0].get("path", "")
    return True, {
        "reason": "app-cli-available",
        "app_path": path,
        "cli_probe_errors": cli_probe_errors,
    }


def _attempt_pencil_cli_execution(
    task: dict[str, Any],
    payload: dict[str, Any],
    dependency_inputs: list[dict[str, Any]],
    gateway_meta: dict[str, Any],
) -> dict[str, Any]:
    cli_command = [str(item) for item in gateway_meta.get("cli_command", []) if str(item).strip()]
    if not cli_command:
        return {
            "status": "skipped",
            "reason": "missing-cli-command",
        }

    script_parts = [
        "Generate UI deliverables from this production orchestration context.",
        f"Package: {task.get('package', '')}",
        f"Intent: {task.get('intent', '')}",
        f"Goal: {payload.get('goal', '')}",
        f"Required files: {', '.join(str(item.get('name', '')) for item in task.get('delivery_contract', {}).get('required_files', []))}",
        f"Dependency inputs: {len(dependency_inputs)}",
    ]
    command = [*cli_command, "-c", "\n".join(script_parts)]
    return_code, stdout, stderr = run_cmd(command, timeout=120)
    if return_code != 0:
        return {
            "status": "failed",
            "reason": "cli-execution-failed",
            "return_code": return_code,
            "command": command,
            "stderr_preview": (stderr or "")[:500],
            "stdout_preview": (stdout or "")[:240],
        }
    return {
        "status": "success",
        "reason": "cli-execution-ok",
        "return_code": return_code,
        "command": command,
        "stdout_preview": (stdout or "")[:500],
    }


def _resolve_gateway(
    selected_executor: str,
    selected_candidate: dict[str, Any],
    discovery_report: dict[str, Any],
    gateway_mode: str,
) -> tuple[bool, str, list[dict[str, Any]], dict[str, Any]]:
    trace: list[dict[str, Any]] = []
    last_meta: dict[str, Any] = {}

    if not selected_executor:
        trace.append({"gateway": "none", "status": "failed", "reason": "no-executor-selected"})
        return False, "", trace, {"reason": "no-executor-selected"}

    if selected_candidate.get("plugin_available") and selected_executor in {"figma_connector", "canva"}:
        trace.append(
            {
                "gateway": "direct",
                "status": "success",
                "reason": "plugin-connector-enabled",
            }
        )
        return True, "direct", trace, {"reason": "plugin-connector-enabled"}

    if gateway_mode == "direct":
        sequence = ["direct"]
    elif gateway_mode == "swarm":
        sequence = ["swarm"]
    else:
        sequence = ["direct", "swarm", "app-cli"]

    for gateway in sequence:
        if gateway == "direct":
            if not selected_candidate.get("direct_available"):
                trace.append({"gateway": "direct", "status": "skipped", "reason": "direct-not-available"})
                continue
            ok, meta = _attempt_direct_gateway(selected_candidate)
            trace.append(
                {
                    "gateway": "direct",
                    "status": "success" if ok else "failed",
                    "reason": meta.get("reason", ""),
                }
            )
            last_meta = meta
            if ok:
                return True, "direct", trace, meta
        elif gateway == "swarm":
            if not selected_candidate.get("swarm_available"):
                trace.append({"gateway": "swarm", "status": "skipped", "reason": "swarm-not-available"})
                continue
            ok, meta = _attempt_swarm_gateway(discovery_report)
            trace.append(
                {
                    "gateway": "swarm",
                    "status": "success" if ok else "failed",
                    "reason": meta.get("reason", ""),
                }
            )
            last_meta = meta
            if ok:
                return True, "swarm", trace, meta
        elif gateway == "app-cli":
            if not selected_candidate.get("app_cli_available"):
                trace.append({"gateway": "app-cli", "status": "skipped", "reason": "app-cli-not-available"})
                continue
            ok, meta = _attempt_app_cli_gateway(selected_candidate)
            trace.append(
                {
                    "gateway": "app-cli",
                    "status": "success" if ok else "failed",
                    "reason": meta.get("reason", ""),
                }
            )
            last_meta = meta
            if ok:
                return True, "app-cli", trace, meta

    if not trace:
        trace.append({"gateway": gateway_mode, "status": "failed", "reason": "no-gateway-attempted"})
    return False, "", trace, last_meta


def _simple_svg(width: int, height: int, title: str, subtitle: str, accent: str) -> str:
    return f"""<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">
  <defs>
    <linearGradient id="g" x1="0" x2="1" y1="0" y2="1">
      <stop offset="0%" stop-color="{accent}" stop-opacity="0.95" />
      <stop offset="100%" stop-color="#111827" stop-opacity="0.95" />
    </linearGradient>
  </defs>
  <rect x="0" y="0" width="{width}" height="{height}" fill="#f8fafc" />
  <circle cx="{int(width * 0.2)}" cy="{int(height * 0.5)}" r="{int(min(width, height) * 0.16)}" fill="url(#g)" />
  <rect x="{int(width * 0.35)}" y="{int(height * 0.36)}" width="{int(width * 0.48)}" height="{int(height * 0.28)}" rx="{int(min(width, height) * 0.04)}" fill="#0f172a" />
  <text x="{int(width * 0.39)}" y="{int(height * 0.49)}" font-family="Helvetica, Arial, sans-serif" font-size="{int(height * 0.08)}" fill="#e2e8f0">{title}</text>
  <text x="{int(width * 0.39)}" y="{int(height * 0.58)}" font-family="Helvetica, Arial, sans-serif" font-size="{int(height * 0.04)}" fill="#94a3b8">{subtitle}</text>
</svg>
"""


def _write_pdf(path: Path, title: str, lines: list[str]) -> None:
    if HAS_REPORTLAB:
        c = canvas.Canvas(str(path), pagesize=A4)
        width, height = A4
        c.setFont("Helvetica-Bold", 16)
        c.drawString(48, height - 72, title)
        c.setFont("Helvetica", 10)
        y = height - 96
        for line in lines:
            if y < 56:
                c.showPage()
                c.setFont("Helvetica", 10)
                y = height - 56
            c.drawString(48, y, line[:120])
            y -= 14
        c.save()
        return

    escaped = [line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)") for line in lines]
    stream_lines = [
        "BT",
        "/F1 14 Tf",
        "48 780 Td",
        f"({title.replace('(', '[').replace(')', ']')}) Tj",
        "/F1 10 Tf",
    ]
    y = 760
    for line in escaped[:40]:
        stream_lines.append(f"1 0 0 1 48 {y} Tm ({line}) Tj")
        y -= 14
    stream_lines.append("ET")
    stream = "\n".join(stream_lines).encode("latin-1", errors="replace")
    objects = [
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n",
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n",
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >> endobj\n",
        b"4 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n",
        f"5 0 obj << /Length {len(stream)} >> stream\n".encode("latin-1") + stream + b"\nendstream endobj\n",
    ]
    offset = 0
    chunks = [b"%PDF-1.4\n"]
    offset += len(chunks[0])
    xref = [0]
    for obj in objects:
        xref.append(offset)
        chunks.append(obj)
        offset += len(obj)
    xref_start = offset
    chunks.append(f"xref\n0 {len(xref)}\n".encode("latin-1"))
    chunks.append(b"0000000000 65535 f \n")
    for item in xref[1:]:
        chunks.append(f"{item:010d} 00000 n \n".encode("latin-1"))
    chunks.append(
        f"trailer << /Size {len(xref)} /Root 1 0 R >>\nstartxref\n{xref_start}\n%%EOF\n".encode("latin-1")
    )
    path.write_bytes(b"".join(chunks))


def _write_ai(path: Path, title: str) -> None:
    content = f"""%!PS-Adobe-3.0
%%Creator: VIS STUDIO 2.2.1
%%Title: {title}
%%BoundingBox: 0 0 800 600
%%Pages: 1
%%EndComments
/Helvetica findfont 36 scalefont setfont
0.12 0.22 0.52 setrgbcolor
120 360 moveto ({title}) show
0.95 0.96 0.98 setrgbcolor
newpath 80 260 moveto 720 260 lineto 720 180 lineto 80 180 lineto closepath fill
0.10 0.16 0.24 setrgbcolor
/Helvetica findfont 18 scalefont setfont
100 220 moveto (Editable Adobe Illustrator source prepared by VIS STUDIO.) show
showpage
%%EOF
"""
    _write_text(path, content)


def _write_pptx(path: Path, title: str, bullets: list[str]) -> None:
    if not HAS_PPTX:
        _write_text(path.with_suffix(".md"), "\n".join([f"# {title}", "", *[f"- {item}" for item in bullets]]))
        return
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    first = True
    for item in bullets:
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = item
    prs.save(str(path))


def _materialize_delivery_files(
    task: dict[str, Any],
    outputs_dir: Path,
    payload: dict[str, Any],
    order_context: dict[str, Any],
    selected_executor: str,
    selected_gateway: str,
    dependency_inputs: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], list[str]]:
    outputs_dir.mkdir(parents=True, exist_ok=True)
    package = task.get("package", "deliverable")
    brand_slug = slugify(str(payload.get("goal", "brand")))
    delivery_contract = task.get("delivery_contract", {})
    required_files = delivery_contract.get("required_files", [])

    created: list[dict[str, Any]] = []
    errors: list[str] = []

    for spec in required_files:
        name = str(spec.get("name", "")).strip()
        if not name:
            continue
        artifact_type = str(spec.get("artifact_type", "artifact"))
        fmt = str(spec.get("format", "")).lower()
        editable = bool(spec.get("editable", True))
        path = outputs_dir / name

        try:
            if fmt == "svg":
                if "poster_sample" in name:
                    content = _simple_svg(
                        1080,
                        1350,
                        "Brand Poster",
                        "Application Sample",
                        "#0b5fff",
                    )
                elif "pattern" in name:
                    content = _simple_svg(1200, 800, "Pattern Library", brand_slug.upper(), "#0891b2")
                elif "brand_graphics" in name:
                    content = _simple_svg(1200, 800, "Brand Graphics", brand_slug.upper(), "#2563eb")
                else:
                    content = _simple_svg(1200, 600, "Brand Logo", brand_slug.upper(), "#1d4ed8")
                _write_text(path, content)
            elif fmt == "md":
                if name == "VIIS_MASTER.md":
                    content = [
                        "# Brand Identity System (VIIS)",
                        "",
                        f"- Package: `{package}`",
                        f"- Executor: `{selected_executor}`",
                        f"- Gateway: `{selected_gateway}`",
                        "",
                        "## Basic Specification",
                        "- Logo system, lockups, safe area, minimum size, misuse.",
                        "- Standard colors (primary/secondary/neutral and accessibility contrast).",
                        "- Typography hierarchy and tone.",
                        "",
                        "## Application Specification",
                        "- Social template structure and content modules.",
                        "- Poster and campaign adaptation rules.",
                        "- File naming, export settings, and handoff quality checks.",
                    ]
                    _write_text(path, "\n".join(content))
                elif name.endswith("logo_usage_rules.md"):
                    content = [
                        "# Logo Usage Rules",
                        "",
                        "## Allowed",
                        "- Use supplied SVG as source of truth.",
                        "- Keep clear space >= 0.5x logo height.",
                        "",
                        "## Not Allowed",
                        "- Distortion, stroke replacement, gradient remapping, random shadows.",
                        "",
                        "## Export",
                        "- Digital: SVG, PNG.",
                        "- Print: AI/PDF vector source.",
                    ]
                    _write_text(path, "\n".join(content))
                elif name.startswith("PROMPTS_"):
                    platform = name.replace("PROMPTS_", "").replace(".md", "")
                    content = [
                        f"# {platform} Prompt Pack",
                        "",
                        "Use this with the approved VIIS constraints.",
                        "",
                        "## Base Prompt",
                        "Create a campaign visual following the approved brand identity system.",
                        "",
                        "## Inject Variables",
                        "- `{headline}`",
                        "- `{subheadline}`",
                        "- `{cta}`",
                        "- `{layout_mode}`",
                        "- `{color_variant}`",
                    ]
                    _write_text(path, "\n".join(content))
                else:
                    _write_text(path, f"# {name}\n\nGenerated by VIS STUDIO 2.2.1.\n")
            elif fmt == "json":
                if name == "PROMPT_VARIABLES.json":
                    data = {
                        "headline": "",
                        "subheadline": "",
                        "cta": "",
                        "layout_mode": "portrait",
                        "color_variant": "primary",
                    }
                else:
                    data = {
                        "package": package,
                        "generated_at": datetime.now(timezone.utc).isoformat(),
                        "selected_executor": selected_executor,
                        "selected_gateway": selected_gateway,
                        "dependency_count": len(dependency_inputs),
                    }
                _write_json(path, data)
            elif fmt == "pdf":
                _write_pdf(
                    path,
                    title=name.replace(".pdf", ""),
                    lines=[
                        f"Package: {package}",
                        f"Executor: {selected_executor}",
                        f"Gateway: {selected_gateway}",
                        "This document is part of a sign-off-ready Brand Identity System delivery.",
                    ],
                )
            elif fmt == "ai":
                _write_ai(path, title=name.replace(".ai", ""))
            elif fmt == "pptx":
                _write_pptx(
                    path,
                    title=name.replace(".pptx", ""),
                    bullets=[
                        "Brand system overview",
                        "Basic identity standards",
                        "Application standards",
                        "Usage checklist",
                    ],
                )
            else:
                _write_text(path, f"Generated file for {name}.\n")
        except Exception as exc:  # defensive: one file failing should not crash full run immediately
            errors.append(f"{name}: {exc}")
            continue

        actual_path = path
        if fmt == "pptx" and not path.exists() and path.with_suffix(".md").exists():
            actual_path = path.with_suffix(".md")

        if not actual_path.exists() or actual_path.stat().st_size == 0:
            errors.append(f"{name}: file missing or empty after generation")
            continue

        created.append(
            {
                "name": actual_path.name,
                "artifact_type": artifact_type,
                "format": actual_path.suffix.replace(".", "").lower(),
                "path": str(actual_path.resolve()),
                "editable": editable,
                "source": {
                    "intent": task.get("intent", ""),
                    "executor": selected_executor,
                    "gateway": selected_gateway,
                },
                "dependency_inputs": dependency_inputs,
            }
        )
    return created, errors


def _execute_task(
    task: dict[str, Any],
    selected_executor: str,
    selected_candidate: dict[str, Any],
    discovery_report: dict[str, Any],
    executor_mode: str,
    execution_gateway_mode: str,
    family_output_dir: Path,
    payload: dict[str, Any],
    order_context: dict[str, Any],
    dependency_inputs: list[dict[str, Any]],
) -> tuple[bool, dict[str, Any], list[dict[str, Any]]]:
    meta: dict[str, Any] = {
        "executor_mode": executor_mode,
        "selected_executor": selected_executor,
        "timestamp": datetime.now(timezone.utc).isoformat(),
    }
    if not selected_executor:
        meta["error"] = "no-executor-available"
        return False, meta, []

    gateway_ok, gateway, gateway_trace, gateway_meta = _resolve_gateway(
        selected_executor=selected_executor,
        selected_candidate=selected_candidate,
        discovery_report=discovery_report,
        gateway_mode=execution_gateway_mode,
    )
    meta["gateway_trace"] = gateway_trace
    meta["gateway_meta"] = gateway_meta
    if not gateway_ok:
        meta["error"] = "gateway-unavailable"
        return False, meta, []

    if selected_executor == "pencil" and gateway == "app-cli":
        pencil_cli = _attempt_pencil_cli_execution(
            task=task,
            payload=payload,
            dependency_inputs=dependency_inputs,
            gateway_meta=gateway_meta,
        )
        meta["pencil_cli_execution"] = pencil_cli

    outputs, errors = _materialize_delivery_files(
        task=task,
        outputs_dir=family_output_dir,
        payload=payload,
        order_context=order_context,
        selected_executor=selected_executor,
        selected_gateway=gateway,
        dependency_inputs=dependency_inputs,
    )
    meta["selected_gateway"] = gateway
    meta["file_count"] = len(outputs)
    if errors:
        meta["generation_errors"] = errors
    success = len(outputs) > 0 and not errors
    meta["status"] = "success" if success else "failed"
    return success, meta, outputs


def _run_document_mode(args: argparse.Namespace) -> int:
    output_root = args.output_root.resolve()
    output_root.mkdir(parents=True, exist_ok=True)

    plan_ns = argparse.Namespace(
        goal=args.goal,
        intent=args.intent,
        execution_mode=args.execution_mode,
        route_rules=args.route_rules,
        preferred_tool=args.preferred_tool,
        template_file=None,
        prompt_template_dir=None,
        output=None,
    )
    payload_ns = argparse.Namespace(
        goal=args.goal,
        intent=args.intent,
        execution_mode=args.execution_mode,
        run_mode="document",
        production_scope="single",
        execution_gateway_mode=args.execution_gateway_mode,
        decision_timeout_seconds=max(1, args.decision_timeout_seconds),
        fallback_figma_connector=bool(args.fallback_figma_connector),
        route_rules=args.route_rules,
        preferred_tool=args.preferred_tool,
        template_dir=None,
        target_file=None,
        target_selection=None,
        selected_executor=None,
        output_family_dir="",
        dependency_inputs_file=None,
        fallback_trace_file=None,
        order_context_file=None,
        gateway_trace_file=None,
        delivery_contract_file=None,
        output=None,
    )

    plan_text = build_plan(plan_ns)
    payload = build_payload(payload_ns)

    _write_text(output_root / "EXECUTION_PLAN.md", plan_text)
    _write_json(output_root / "TOOL_EXECUTION_PAYLOAD.json", payload)
    print(f"Wrote: {output_root / 'EXECUTION_PLAN.md'}")
    print(f"Wrote: {output_root / 'TOOL_EXECUTION_PAYLOAD.json'}")
    return 0


def _run_production_mode(args: argparse.Namespace, routes: dict[str, Any]) -> int:
    if args.brand_pack is not None and args.brand_pack.exists():
        try:
            pack = json.loads(args.brand_pack.read_text(encoding="utf-8"))
        except Exception:
            pack = {}
        brand_name = (
            pack.get("brand", {}).get("slug")
            or pack.get("brand", {}).get("name")
            or "brand"
        )
    else:
        brand_name = "brand"
    brand_slug = slugify(str(brand_name))

    run_id = f"{utc_timestamp()}_{brand_slug}"
    run_root = args.output_root.resolve() / "production_runs" / run_id
    run_root.mkdir(parents=True, exist_ok=True)

    run_log: list[str] = []

    def log(message: str) -> None:
        now = datetime.now(timezone.utc).strftime("%H:%M:%S")
        line = f"- [{now} UTC] {message}"
        run_log.append(line)
        print(message, flush=True)

    order_context, order_file = _load_or_create_order_form(
        order_file=args.order_file,
        run_root=run_root,
        resume_override=args.resume_from_run_dir,
    )
    log(f"Order form prepared: {order_file}")

    discovery_report = discover_design_services(routes)
    _write_json(run_root / "MCP_DISCOVERY_REPORT.json", discovery_report)
    if discovery_report.get("missing_critical"):
        _build_setup_required_doc(run_root / "MCP_SETUP_REQUIRED.md", discovery_report)
    log("MCP discovery completed.")

    pipeline_ns = argparse.Namespace(
        goal=args.goal,
        intent=args.intent,
        production_scope=args.production_scope,
        brand_pack=args.brand_pack,
        order_file=order_file,
        resume_from_run_dir=args.resume_from_run_dir,
        preferred_tool=args.preferred_tool,
        route_rules=args.route_rules,
        output=None,
    )
    pipeline = build_pipeline(pipeline_ns)
    order_context = pipeline.get("order_context", order_context)
    _write_json(run_root / "ORDER_FORM.json", order_context)
    _write_json(run_root / "PIPELINE.json", pipeline)
    log(f"Pipeline built with {pipeline.get('task_count', 0)} tasks.")

    artifact_registry: dict[str, Any] = {
        "run_id": run_id,
        "order_context": order_context,
        "packages": {},
        "tasks": {},
        "files": [],
        "by_artifact_type": {},
    }
    manifest: dict[str, Any] = {
        "run_id": run_id,
        "goal": args.goal,
        "run_mode": "production",
        "production_scope": args.production_scope,
        "execution_mode": "execute",
        "execution_gateway_mode": args.execution_gateway_mode,
        "executor_mode": args.executor_mode,
        "fallback_policy": {
            "figma_connector": bool(args.fallback_figma_connector),
        },
        "critical_failure_default_action": "stop",
        "decision_timeout_seconds": max(1, args.decision_timeout_seconds),
        "order_file": str(order_file),
        "started_at": datetime.now(timezone.utc).isoformat(),
        "status": "running",
    }

    missing_apps: list[str] = []
    stopped = False
    for task in pipeline.get("tasks", []):
        task_id = task["task_id"]
        family = slugify(task.get("family", "deliverable"))
        family_dir = run_root / family
        for child in ("inputs", "requests", "outputs", "meta"):
            (family_dir / child).mkdir(parents=True, exist_ok=True)

        if task.get("skip_existing"):
            resume_info = task.get("resume_package", {})
            artifact_registry["tasks"][task_id] = {
                "package": task.get("package", ""),
                "intent": task["intent"],
                "family": task.get("family"),
                "status": "skipped",
                "depends_on": task.get("depends_on", []),
                "selected_executor": "",
                "selected_gateway": "resume",
                "output_files": resume_info.get("files", []),
            }
            artifact_registry["packages"][task.get("package", "")] = {
                "status": "skipped",
                "task_id": task_id,
                "files": resume_info.get("files", []),
            }
            log(f"Task skipped from resume: {task_id}")
            continue

        selected_executor, candidates, fallback_trace, selected_reason, selected_candidate = (
            _select_executor_for_task(
                task=task,
                discovery_report=discovery_report,
                fallback_figma_connector=args.fallback_figma_connector,
            )
        )
        dependency_inputs = _collect_dependency_inputs(task, artifact_registry)
        _write_json(family_dir / "inputs" / "dependency_inputs.json", dependency_inputs)

        payload_ns = argparse.Namespace(
            goal=f"{args.goal} [{task.get('package', task.get('family', 'task'))}]",
            intent=task["intent"],
            execution_mode="execute",
            run_mode="production",
            production_scope=args.production_scope,
            execution_gateway_mode=args.execution_gateway_mode,
            decision_timeout_seconds=max(1, args.decision_timeout_seconds),
            fallback_figma_connector=bool(args.fallback_figma_connector),
            route_rules=args.route_rules,
            preferred_tool=args.preferred_tool,
            template_dir=None,
            target_file=None,
            target_selection=None,
            selected_executor=selected_executor or None,
            output_family_dir=str(family_dir),
            dependency_inputs_file=None,
            fallback_trace_file=None,
            order_context_file=None,
            gateway_trace_file=None,
            delivery_contract_file=None,
            output=None,
        )
        payload = build_payload(payload_ns)
        payload["executor_candidates"] = task.get("executor_candidates", [])
        payload["selected_executor"] = selected_executor
        payload["dependency_inputs"] = dependency_inputs
        payload["output_family_dir"] = str(family_dir)
        payload["fallback_trace"] = fallback_trace
        payload["selection_reason"] = selected_reason
        payload["depends_on"] = task.get("depends_on", [])
        payload["critical"] = bool(task.get("critical", False))
        payload["order_context"] = order_context
        payload["delivery_contract"] = task.get("delivery_contract", {})
        payload["gateway_trace"] = []
        _write_json(family_dir / "requests" / "TOOL_EXECUTION_PAYLOAD.json", payload)

        plan_ns = argparse.Namespace(
            goal=f"{args.goal} [{task.get('package', task.get('family', 'task'))}]",
            intent=task["intent"],
            execution_mode="execute",
            route_rules=args.route_rules,
            preferred_tool=args.preferred_tool,
            template_file=None,
            prompt_template_dir=None,
            output=None,
        )
        _write_text(family_dir / "requests" / "EXECUTION_PLAN.md", build_plan(plan_ns))

        success, exec_meta, output_files = _execute_task(
            task=task,
            selected_executor=selected_executor,
            selected_candidate=selected_candidate,
            discovery_report=discovery_report,
            executor_mode=args.executor_mode,
            execution_gateway_mode=args.execution_gateway_mode,
            family_output_dir=family_dir / "outputs",
            payload=payload,
            order_context=order_context,
            dependency_inputs=dependency_inputs,
        )
        payload["gateway_trace"] = exec_meta.get("gateway_trace", [])
        payload["selected_gateway"] = exec_meta.get("selected_gateway", "")
        _write_json(family_dir / "requests" / "TOOL_EXECUTION_PAYLOAD.json", payload)

        if not success and task.get("critical", False):
            if exec_meta.get("error") == "gateway-unavailable":
                app_hits = selected_candidate.get("app_cli_hits", [])
                if not app_hits and selected_candidate.get("app_cli_available") is False:
                    app_aliases = selected_candidate.get("app_cli_hits", [])
                    if not app_aliases:
                        executor = selected_candidate.get("executor", selected_executor)
                        if executor:
                            missing_apps.append(executor)
            decision = _prompt_for_critical_decision(task_id, args.decision_timeout_seconds)
            exec_meta["critical_decision"] = decision
            if decision == "retry":
                success, retry_meta, output_files = _execute_task(
                    task=task,
                    selected_executor=selected_executor,
                    selected_candidate=selected_candidate,
                    discovery_report=discovery_report,
                    executor_mode=args.executor_mode,
                    execution_gateway_mode=args.execution_gateway_mode,
                    family_output_dir=family_dir / "outputs",
                    payload=payload,
                    order_context=order_context,
                    dependency_inputs=dependency_inputs,
                )
                exec_meta["retry"] = retry_meta
            elif decision == "skip":
                success = False
                exec_meta["status"] = "critical-skipped-by-user"
            else:
                exec_meta["status"] = "critical-stopped"
                _write_json(family_dir / "meta" / "execution_meta.json", exec_meta)
                artifact_registry["tasks"][task_id] = {
                    "package": task.get("package", ""),
                    "intent": task["intent"],
                    "family": task.get("family"),
                    "status": "stopped",
                    "depends_on": task.get("depends_on", []),
                    "selected_executor": selected_executor,
                    "selected_gateway": exec_meta.get("selected_gateway", ""),
                    "output_files": output_files,
                }
                artifact_registry["packages"][task.get("package", "")] = {
                    "status": "stopped",
                    "task_id": task_id,
                    "files": output_files,
                }
                log(f"Critical task stopped: {task_id}")
                stopped = True
                break

        status = "success" if success else ("failed" if task.get("critical", False) else "failed-noncritical")
        if not success and not task.get("critical", False):
            log(f"Non-critical task failed and continued: {task_id}")
        elif success:
            log(
                f"Task completed: {task_id} via {selected_executor or 'none'} "
                f"[{exec_meta.get('selected_gateway', 'no-gateway')}]"
            )

        _write_json(family_dir / "meta" / "execution_meta.json", exec_meta)
        artifact_registry["tasks"][task_id] = {
            "package": task.get("package", ""),
            "intent": task["intent"],
            "family": task.get("family"),
            "status": status,
            "depends_on": task.get("depends_on", []),
            "selected_executor": selected_executor,
            "selected_gateway": exec_meta.get("selected_gateway", ""),
            "candidates": candidates,
            "output_files": output_files,
        }
        artifact_registry["packages"][task.get("package", "")] = {
            "status": status,
            "task_id": task_id,
            "files": output_files,
        }
        for item in output_files:
            artifact_registry["files"].append(item)
            atype = item.get("artifact_type", "")
            artifact_registry["by_artifact_type"].setdefault(atype, []).append(
                {
                    "task_id": task_id,
                    "package": task.get("package", ""),
                    "path": item.get("path", ""),
                    "format": item.get("format", ""),
                }
            )

    if missing_apps:
        _build_app_cli_setup_doc(run_root / "APP_CLI_SETUP_REQUIRED.md", sorted(set(missing_apps)))

    manifest["finished_at"] = datetime.now(timezone.utc).isoformat()
    task_states = [item.get("status", "") for item in artifact_registry.get("tasks", {}).values()]
    if stopped:
        manifest["status"] = "stopped"
    elif any(state.startswith("failed") for state in task_states):
        manifest["status"] = "completed-with-failures"
    else:
        manifest["status"] = "completed"
    manifest["task_summary"] = {
        "total": len(task_states),
        "success": len([state for state in task_states if state == "success"]),
        "failed": len([state for state in task_states if state.startswith("failed")]),
        "skipped": len([state for state in task_states if state == "skipped"]),
        "stopped": len([state for state in task_states if state == "stopped"]),
    }

    _write_json(run_root / "ARTIFACT_REGISTRY.json", artifact_registry)
    _write_json(run_root / "RUN_MANIFEST.json", manifest)
    _write_text(run_root / "RUN_LOG.md", "# Run Log\n\n" + "\n".join(run_log) + "\n")

    print(f"Production run created: {run_root}")
    print(f"Manifest: {run_root / 'RUN_MANIFEST.json'}")
    return 0 if manifest["status"] in {"completed", "completed-with-failures"} else 2


def main() -> None:
    args = parse_args()
    routes = load_route_rules(args.route_rules)
    if args.run_mode == "document":
        raise_code = _run_document_mode(args)
    else:
        raise_code = _run_production_mode(args, routes)
    raise SystemExit(raise_code)


if __name__ == "__main__":
    main()
